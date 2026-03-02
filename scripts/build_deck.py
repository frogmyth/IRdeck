"""AIsirius PPTX 빌드 스크립트.

마크다운 콘텐츠 → PPTX 자동 생성.

Usage:
    python scripts/build_deck.py --version A
    python scripts/build_deck.py --version B
    python scripts/build_deck.py --version all
    python scripts/build_deck.py --version A --watermark --company "대상회사"
"""
import argparse
import io
import os
import re
import sys
from datetime import date

# Windows cp949 인코딩 문제 해결
if sys.stdout.encoding != "utf-8":
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8", errors="replace")

# 프로젝트 루트를 path에 추가
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, SCRIPT_DIR)

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from components.colors import (
    DARK_NAVY, VERY_DARK_NAVY, BRAND_CYAN, STANDARD_BLUE,
    BRIGHT_BLUE, TEAL, WHITE, BLACK, PALE_BLUE,
    MUTED_TEAL, LIME_GREEN, TABLE_HEADER_BG, TABLE_HEADER_FG,
    TABLE_ALT_ROW, TABLE_BODY_FG,
)
from components.fonts import apply_font
from components.layouts import (
    SLIDE_WIDTH, SLIDE_HEIGHT,
    CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT,
    A_TEXT_LEFT, A_TEXT_WIDTH, A_IMAGE_LEFT, A_IMAGE_WIDTH,
    B_TEXT_LEFT, B_TEXT_WIDTH, B_IMAGE_LEFT, B_IMAGE_WIDTH,
    CHART_3_WIDTH, CHART_3_HEIGHT, CHART_3_POSITIONS,
    CHART_SINGLE_LEFT, CHART_SINGLE_TOP, CHART_SINGLE_WIDTH, CHART_SINGLE_HEIGHT,
    PILLAR_WIDTH, PILLAR_HEIGHT, PILLAR_GAP, PILLAR_POSITIONS,
    LayoutType,
)
from components.chrome import add_full_chrome
from components.markdown_parser import parse_slides, SlideContent
from components.charts import (
    CHART_BUILDERS, add_esl_adoption_table,
    add_source_text, MARKET_DATA,
)

# 프로젝트 경로
PROJECT_ROOT = os.path.dirname(SCRIPT_DIR)
TEMPLATES_DIR = os.path.join(PROJECT_ROOT, "templates")
CONTENT_DIR = os.path.join(PROJECT_ROOT, "docs", "content")
OUTPUT_DIR = os.path.join(PROJECT_ROOT, "output", "회사소개서")
ASSETS_DIR = os.path.join(PROJECT_ROOT, "assets")


# ── 섹션 매핑 (슬라이드 번호 → 섹션 정보) ────────────
SECTION_MAP_A = {
    1: ("", ""),
    2: ("AIsirius 소개", "About AIsirius"),
    3: ("시장 기회", "Marketing Opportunity"),
    4: ("시장 기회", "Expected Outcomes and Target Market"),
    5: ("시장 기회", "Retail Business Strategies"),
    6: ("시장 기회", "Expected Outcomes and Target Market"),
    7: ("AIsirius AI", "Core AI Technology"),
    8: ("AIsirius AI", "Core AI Technology"),
    9: ("AIsirius AI", "AI Content Generation"),
    10: ("AIsirius AI", "AI Store Analysis"),
    11: ("AIsirius AI", "3-Stage Distributed AI"),
    12: ("플랫폼", "Platform Architecture"),
    13: ("플랫폼", "CMS Features"),
    14: ("플랫폼", "Cross-Device Integration"),
    15: ("플랫폼", "HW Lineup"),
    16: ("플랫폼", "Device Absorption"),
    17: ("비즈니스", "Business Model"),
    18: ("비즈니스", "ROI Analysis"),
    19: ("비즈니스", "Global Strategy"),
    20: ("비즈니스", "Revenue Roadmap"),
    21: ("팀", "About AIsirius"),
    22: ("팀", "Traction"),
    23: ("팀", "ESG"),
    24: ("", ""),
    25: ("부록", "Appendix"),
}

SECTION_MAP_B = {
    1: ("", ""),                                             # 표지
    2: ("AIsirius 소개", "Company Definition"),               # 한 줄 정의 (신규)
    3: ("시장 기회", "Marketing Opportunity"),                 # DX→AX
    4: ("시장 기회", "Marketing Opportunity"),                 # AX 비즈니스 기회
    5: ("시장 기회", "Expected Outcomes and Target Market"),   # 글로벌 시장 규모
    6: ("시장 기회", "Expected Outcomes and Target Market"),   # ESL 지역별 도입률
    7: ("시장 기회", "Retail Business Strategies"),            # 월마트 사례 1/2
    8: ("시장 기회", "Retail Business Strategies"),            # 월마트 사례 2/2
    9: ("비즈니스 모델", "Business Model"),                    # AIsirius 핵심 정의
    10: ("비즈니스 모델", "ISO4OM 3 Solutions"),               # ISO4OM 3 Solutions
    11: ("AI 기술", "AI Technology"),                          # 전용 AI vs 범용 AI
    12: ("플랫폼", "CMS Platform"),                           # CMS 기능 1/2
    13: ("플랫폼", "CMS Platform"),                           # CMS 기능 2/2
    14: ("AI 기술", "AI Content Generation"),                  # AI 콘텐츠 자동 생성
    15: ("플랫폼", "Cross-Device Integration"),                # 크로스 디바이스
    16: ("AI 기술", "3-Stage Distributed AI"),                 # 3단 분산 AI
    17: ("플랫폼", "HW Lineup"),                              # HW 라인업
    18: ("비즈니스", "Revenue Model"),                         # 수익 모델
    19: ("비즈니스", "ROI Analysis"),                          # ROI 분석
    20: ("비즈니스", "Global Strategy"),                       # 글로벌 전략
    21: ("비즈니스", "Revenue Roadmap"),                       # J-Curve
    22: ("비즈니스", "Expansion Vision"),                      # 통합 솔루션 확장
    23: ("팀", "Traction"),                                   # 국내 트랙션
    24: ("팀", "ESG"),                                        # ESG
    25: ("팀", "About AIsirius"),                             # CEO & 팀
    26: ("", ""),                                             # 감사합니다
    27: ("부록", "Appendix"),                                 # 부록 1
    28: ("부록", "Appendix"),                                 # 부록 2
    29: ("부록", "Appendix"),                                 # 부록 3
}


# ══════════════════════════════════════════════════════
#  레이아웃 빌더 함수들
# ══════════════════════════════════════════════════════

def _add_text_box(slide, left, top, width, height, text, font_role,
                  color=None, alignment=PP_ALIGN.LEFT, word_wrap=True,
                  override_size=None):
    """텍스트박스를 추가하고 폰트를 적용하는 헬퍼."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    apply_font(run, font_role, color=color, override_size=override_size)
    return txbox


def _add_bullet_list(slide, left, top, width, height, items, font_role,
                     color=None, spacing=Pt(4)):
    """불릿 리스트 텍스트박스 추가."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing
        p.level = 0

        # Bold 마크다운 처리
        parts = re.split(r"(\*\*[^*]+\*\*)", item)
        for part in parts:
            if part.startswith("**") and part.endswith("**"):
                run = p.add_run()
                run.text = part[2:-2]
                apply_font(run, font_role, override_bold=True, color=color)
            elif part:
                run = p.add_run()
                run.text = part
                apply_font(run, font_role, color=color)
    return txbox


def _add_body_paragraphs(slide, left, top, width, height, texts, font_role,
                         color=None, line_spacing=1.15):
    """본문 텍스트 여러 줄 추가 (Bold 마크다운 지원)."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True

    for i, text in enumerate(texts):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(4)

        # Bold 마크다운 처리
        parts = re.split(r"(\*\*[^*]+\*\*)", text)
        for part in parts:
            if part.startswith("**") and part.endswith("**"):
                run = p.add_run()
                run.text = part[2:-2]
                apply_font(run, font_role, override_bold=True, color=color)
            elif part:
                run = p.add_run()
                run.text = part
                apply_font(run, font_role, color=color)
    return txbox


def _add_image_placeholder(slide, left, top, width, height, description=""):
    """이미지 플레이스홀더 — 모던 스타일 (연한 파란 배경 + 아이콘)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PALE_BLUE
    shape.line.color.rgb = RGBColor(0xB8, 0xD4, 0xE8)
    shape.line.width = Pt(0.75)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    # 아이콘 역할의 심볼
    icon_run = tf.paragraphs[0].add_run()
    icon_run.text = "\u25A1"  # □ 사각형 심볼
    apply_font(icon_run, "main_heading", override_size=Pt(28),
               color=RGBColor(0xB8, 0xD4, 0xE8))
    # 설명 텍스트
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    desc_run = p2.add_run()
    desc_run.text = description if description else "Image"
    apply_font(desc_run, "source_text", color=MUTED_TEAL)
    return shape


def _add_styled_table(slide, left, top, width, height, table_data):
    """파싱된 테이블 데이터로 스타일 테이블 추가."""
    if not table_data or not table_data.get("headers"):
        return None

    headers = table_data["headers"]
    rows = table_data["rows"]
    n_rows = len(rows) + 1
    n_cols = len(headers)

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # 헤더 행
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(10)
            para.font.bold = True
            para.font.color.rgb = TABLE_HEADER_FG
            para.alignment = PP_ALIGN.CENTER

    # 데이터 행
    for i, row in enumerate(rows):
        for j, val in enumerate(row[:n_cols]):
            cell = table.cell(i + 1, j)
            # Bold 마크다운 제거
            clean_val = re.sub(r"\*\*([^*]+)\*\*", r"\1", val)
            cell.text = clean_val
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_ROW
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(9)
                para.font.color.rgb = TABLE_BODY_FG
                para.alignment = PP_ALIGN.CENTER

    return table_shape


# ── 슬라이드 타입별 빌더 ─────────────────────────────

def build_title_cover(slide, sc, version):
    """표지 슬라이드 — 모던 디자인 (크롬 없음)."""
    # 배경 — 네이비 풀 블리드
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_NAVY
    bg.line.fill.background()

    # 좌측 Cyan 액센트 스트립
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        Inches(0.12), SLIDE_HEIGHT,
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = BRAND_CYAN
    accent.line.fill.background()

    # 상단 Cyan 수평선 (장식)
    top_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2.0), Inches(1.8),
        Inches(9.0), Inches(0.03),
    )
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = BRAND_CYAN
    top_line.line.fill.background()

    # 로고 텍스트 (큰 사이즈)
    logo_box = slide.shapes.add_textbox(
        Inches(2.0), Inches(2.1), Inches(9.0), Inches(1.2)
    )
    tf = logo_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run_ai = p.add_run()
    run_ai.text = "AI"
    apply_font(run_ai, "main_heading", override_size=Pt(60), color=BRAND_CYAN)
    run_s = p.add_run()
    run_s.text = "sirius"
    apply_font(run_s, "main_heading", override_size=Pt(60), color=WHITE)

    # 서브타이틀
    _add_text_box(
        slide, Inches(2.0), Inches(3.4), Inches(9.0), Inches(0.5),
        "ISO4OM Platform  —  In Store O4O Merchandising",
        "subtitle", color=BRIGHT_BLUE, alignment=PP_ALIGN.LEFT,
    )

    # 메인 태그라인
    _add_text_box(
        slide, Inches(2.0), Inches(4.1), Inches(9.0), Inches(0.7),
        "Store into Media, Shelf into Profit",
        "section_header", color=WHITE, alignment=PP_ALIGN.LEFT,
    )

    # 슬로건
    _add_text_box(
        slide, Inches(2.0), Inches(4.8), Inches(9.0), Inches(0.5),
        "Create AI Smart Flow & Data Driven",
        "header_english", override_size=Pt(16), color=MUTED_TEAL,
        alignment=PP_ALIGN.LEFT,
    )

    # 하단 Cyan 수평선 (장식)
    bot_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2.0), Inches(5.8),
        Inches(9.0), Inches(0.03),
    )
    bot_line.fill.solid()
    bot_line.fill.fore_color.rgb = BRAND_CYAN
    bot_line.line.fill.background()

    # 하단 회사명
    _add_text_box(
        slide, Inches(2.0), Inches(6.0), Inches(5.0), Inches(0.4),
        "에이아이시리우스(주)  |  AIsirius Co., Ltd.",
        "small_label", color=WHITE, alignment=PP_ALIGN.LEFT,
    )

    # 날짜 + Confidential (우측 정렬)
    _add_text_box(
        slide, Inches(8.0), Inches(6.0), Inches(3.0), Inches(0.4),
        f"{date.today().strftime('%Y.%m')}  |  Confidential",
        "source_text", color=MUTED_TEAL, alignment=PP_ALIGN.RIGHT,
    )

    # 하단 3색 액센트 라인
    bar_h = Inches(0.06)
    bar_top = SLIDE_HEIGHT - bar_h
    bar_w = SLIDE_WIDTH / 3
    for i, color in enumerate([DARK_NAVY, STANDARD_BLUE, BRAND_CYAN]):
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(i * 13.333 / 3), bar_top, bar_w, bar_h,
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()


def build_two_column(slide, sc, version):
    """좌우 분할 레이아웃 (텍스트 + 이미지)."""
    if version == "A":
        t_left, t_width = A_TEXT_LEFT, A_TEXT_WIDTH
        i_left, i_width = A_IMAGE_LEFT, A_IMAGE_WIDTH
        font_role = "body_a"
    else:
        t_left, t_width = B_TEXT_LEFT, B_TEXT_WIDTH
        i_left, i_width = B_IMAGE_LEFT, B_IMAGE_WIDTH
        font_role = "body_b"

    # 헤드라인
    if sc.headline:
        _add_text_box(
            slide, t_left, CONTENT_TOP, t_width, Inches(0.6),
            sc.headline, "subtitle", color=DARK_NAVY,
        )
        body_top = CONTENT_TOP + Inches(0.8)
    else:
        body_top = CONTENT_TOP

    # 본문 (불릿 or 텍스트)
    body_height = CONTENT_HEIGHT - (body_top - CONTENT_TOP) - Inches(0.3)
    if sc.bullets:
        _add_bullet_list(
            slide, t_left, body_top, t_width, body_height,
            sc.bullets, font_role, color=VERY_DARK_NAVY,
        )
    elif sc.body_texts:
        _add_body_paragraphs(
            slide, t_left, body_top, t_width, body_height,
            sc.body_texts, font_role, color=VERY_DARK_NAVY,
        )

    # 인용문
    if sc.quotes:
        q_top = body_top + body_height - Inches(1.0)
        for quote in sc.quotes[:1]:
            _add_text_box(
                slide, t_left, q_top, t_width, Inches(0.8),
                f'"{quote}"', "quote_b" if version == "B" else "small_label",
                color=MUTED_TEAL,
            )

    # 이미지 플레이스홀더
    desc = sc.image_refs[0] if sc.image_refs else sc.title
    _add_image_placeholder(
        slide, i_left, CONTENT_TOP, i_width, CONTENT_HEIGHT, desc
    )


def build_chart_slide(slide, sc, version):
    """차트 슬라이드 (시장 데이터)."""
    chart_keys = sc.chart_keys

    if len(chart_keys) >= 3:
        # 3개 차트 병렬 배치
        for i, key in enumerate(chart_keys[:3]):
            if key in CHART_BUILDERS:
                left, top = CHART_3_POSITIONS[i]
                CHART_BUILDERS[key](slide, left, top, CHART_3_WIDTH, CHART_3_HEIGHT)
    elif len(chart_keys) >= 1:
        # 단일 차트
        key = chart_keys[0]
        if key in CHART_BUILDERS:
            CHART_BUILDERS[key](
                slide, CHART_SINGLE_LEFT, CHART_SINGLE_TOP,
                CHART_SINGLE_WIDTH, CHART_SINGLE_HEIGHT
            )

    # 헤드라인 (차트 위에 — 없으면 제목 사용)
    headline = sc.headline or sc.title
    _add_text_box(
        slide, CONTENT_LEFT, CONTENT_TOP - Inches(0.15), CONTENT_WIDTH, Inches(0.5),
        headline, "subtitle", color=DARK_NAVY,
    )


def build_chart_detailed(slide, sc, version):
    """차트 + 상세 텍스트 (Version B)."""
    chart_keys = sc.chart_keys

    # 차트 영역 (좌측 60%)
    chart_width = Inches(7.0)
    chart_height = Inches(4.5)

    if chart_keys:
        key = chart_keys[0]
        if key in CHART_BUILDERS:
            CHART_BUILDERS[key](
                slide, CONTENT_LEFT, CONTENT_TOP + Inches(0.5),
                chart_width, chart_height,
            )

    # 텍스트 영역 (우측 또는 하단)
    text_left = CONTENT_LEFT + chart_width + Inches(0.3)
    text_width = CONTENT_WIDTH - chart_width - Inches(0.3)

    if sc.body_texts:
        _add_body_paragraphs(
            slide, text_left, CONTENT_TOP + Inches(0.5),
            text_width, chart_height,
            sc.body_texts[:10], "body_b", color=VERY_DARK_NAVY,
        )

    # 테이블이 있으면 추가
    if sc.tables:
        table_top = CONTENT_TOP + Inches(0.5) + chart_height + Inches(0.2)
        remaining = CONTENT_HEIGHT - chart_height - Inches(0.7)
        if remaining > Inches(0.5):
            _add_styled_table(
                slide, CONTENT_LEFT, table_top,
                CONTENT_WIDTH, remaining, sc.tables[0],
            )


def build_text_heavy(slide, sc, version):
    """텍스트 중심 슬라이드 (Version B 주력)."""
    top = CONTENT_TOP

    # 헤드라인 (있으면)
    if sc.headline:
        _add_text_box(
            slide, CONTENT_LEFT, top, CONTENT_WIDTH, Inches(0.5),
            sc.headline, "subtitle", color=DARK_NAVY,
        )
        top += Inches(0.7)

    remaining = CONTENT_HEIGHT - (top - CONTENT_TOP)

    # 본문 텍스트 (불릿 + 일반 텍스트 통합)
    all_texts = []
    if sc.body_texts:
        all_texts.extend(sc.body_texts)
    if sc.bullets:
        all_texts.extend([f"• {b}" for b in sc.bullets])

    if all_texts:
        _add_body_paragraphs(
            slide, CONTENT_LEFT, top,
            B_TEXT_WIDTH if version == "B" else CONTENT_WIDTH,
            remaining, all_texts, "body_b" if version == "B" else "body_a",
            color=VERY_DARK_NAVY,
        )

    # 인용문 (있으면)
    if sc.quotes:
        q_top = top + remaining - Inches(1.2)
        for quote in sc.quotes[:2]:
            txbox = _add_text_box(
                slide, CONTENT_LEFT + Inches(0.3), q_top,
                CONTENT_WIDTH - Inches(0.6), Inches(0.5),
                f'> {quote}', "quote_b" if version == "B" else "small_label",
                color=MUTED_TEAL,
            )
            q_top += Inches(0.55)

    # 테이블 (있으면, 우측에)
    if sc.tables and version == "B":
        _add_styled_table(
            slide, B_IMAGE_LEFT, top,
            B_IMAGE_WIDTH, remaining, sc.tables[0],
        )


def build_comparison_table(slide, sc, version):
    """비교 테이블 슬라이드."""
    # 헤드라인
    _add_text_box(
        slide, CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, Inches(0.5),
        sc.headline or sc.title, "subtitle", color=DARK_NAVY,
    )

    # 테이블
    if sc.tables:
        table_top = CONTENT_TOP + Inches(0.7)
        table_height = CONTENT_HEIGHT - Inches(1.0)
        _add_styled_table(
            slide, CONTENT_LEFT, table_top,
            CONTENT_WIDTH, table_height, sc.tables[0],
        )
    # 테이블 없으면 본문 텍스트로 대체
    elif sc.body_texts:
        _add_body_paragraphs(
            slide, CONTENT_LEFT, CONTENT_TOP + Inches(0.7),
            CONTENT_WIDTH, CONTENT_HEIGHT - Inches(1.0),
            sc.body_texts, "body_b" if version == "B" else "body_a",
            color=VERY_DARK_NAVY,
        )


def build_three_pillar(slide, sc, version):
    """3단 (Three Pillar) 레이아웃 — 모던 카드 스타일."""
    # 헤드라인
    _add_text_box(
        slide, CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, Inches(0.5),
        sc.headline or sc.title, "subtitle", color=DARK_NAVY,
    )

    pillar_top = CONTENT_TOP + Inches(0.7)

    # 3개 기둥 (불릿을 3등분)
    items_per_pillar = max(1, len(sc.bullets) // 3)
    pillar_colors = [STANDARD_BLUE, TEAL, LIME_GREEN]

    for i in range(3):
        left = PILLAR_POSITIONS[i]
        start = i * items_per_pillar
        end = start + items_per_pillar if i < 2 else len(sc.bullets)
        pillar_items = sc.bullets[start:end]

        # 상단 컬러 액센트 바 (카드 위)
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, pillar_top,
            PILLAR_WIDTH, Inches(0.08),
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = pillar_colors[i]
        accent_bar.line.fill.background()

        # 기둥 배경 (흰색 카드 + 그림자 효과 = 연한 테두리)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, pillar_top + Inches(0.08),
            PILLAR_WIDTH, PILLAR_HEIGHT - Inches(0.08),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xFA, 0xFB, 0xFC)
        shape.line.color.rgb = RGBColor(0xE0, 0xE4, 0xE8)
        shape.line.width = Pt(0.5)

        # 기둥 제목
        if pillar_items:
            title_text = pillar_items[0] if pillar_items else f"Pillar {i+1}"
            _add_text_box(
                slide, left + Inches(0.25), pillar_top + Inches(0.3),
                PILLAR_WIDTH - Inches(0.5), Inches(0.5),
                title_text, "body_a_bold" if version == "A" else "body_b",
                color=pillar_colors[i], alignment=PP_ALIGN.LEFT,
            )

        # 컬러 구분선
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left + Inches(0.25), pillar_top + Inches(0.85),
            Inches(0.8), Inches(0.02),
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = pillar_colors[i]
        divider.line.fill.background()

        # 기둥 내용
        if len(pillar_items) > 1:
            _add_bullet_list(
                slide, left + Inches(0.25), pillar_top + Inches(1.05),
                PILLAR_WIDTH - Inches(0.5), PILLAR_HEIGHT - Inches(1.3),
                pillar_items[1:],
                "body_a" if version == "A" else "body_b",
                color=VERY_DARK_NAVY,
            )


def build_news_quote(slide, sc, version):
    """뉴스 인용 슬라이드 — 모던 좌측 액센트 인용 스타일."""
    top = CONTENT_TOP

    # 뉴스 인용문들
    for i, quote in enumerate(sc.quotes):
        # 좌측 액센트 바
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            CONTENT_LEFT, top,
            Inches(0.08), Inches(1.1),
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = BRAND_CYAN
        accent.line.fill.background()

        # 인용 박스 배경
        q_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            CONTENT_LEFT + Inches(0.15), top,
            CONTENT_WIDTH - Inches(0.15), Inches(1.1),
        )
        q_box.fill.solid()
        q_box.fill.fore_color.rgb = RGBColor(0xF5, 0xF8, 0xFC)
        q_box.line.fill.background()

        tf = q_box.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].space_before = Pt(8)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = quote
        apply_font(run, "body_b", color=DARK_NAVY)
        top += Inches(1.3)

        if top > Inches(5.5):
            break

    # 나머지 본문
    remaining = CONTENT_HEIGHT - (top - CONTENT_TOP)
    if sc.body_texts and remaining > Inches(0.5):
        _add_body_paragraphs(
            slide, CONTENT_LEFT, top,
            CONTENT_WIDTH, remaining,
            sc.body_texts, "body_b", color=VERY_DARK_NAVY,
        )


def build_infographic_numbers(slide, sc, version):
    """숫자 인포그래픽 — 모던 카드 스타일 (ROI 등)."""
    _add_text_box(
        slide, CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, Inches(0.5),
        sc.headline or sc.title, "subtitle", color=DARK_NAVY,
    )

    # 핵심 숫자를 카드형으로 배치
    numbers_top = CONTENT_TOP + Inches(0.8)
    key_numbers = sc.bullets[:4] if sc.bullets else []
    n_cards = max(1, len(key_numbers))
    col_width = CONTENT_WIDTH / n_cards
    card_colors = [STANDARD_BLUE, TEAL, LIME_GREEN, BRIGHT_BLUE]

    for i, item in enumerate(key_numbers):
        left = CONTENT_LEFT + col_width * i

        # 상단 컬러 액센트
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left + Inches(0.15), numbers_top,
            col_width - Inches(0.3), Inches(0.06),
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = card_colors[i % 4]
        accent.line.fill.background()

        # 카드 배경
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left + Inches(0.15), numbers_top + Inches(0.06),
            col_width - Inches(0.3), Inches(2.2),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xFA, 0xFB, 0xFC)
        shape.line.color.rgb = RGBColor(0xE0, 0xE4, 0xE8)
        shape.line.width = Pt(0.5)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        # 수치/텍스트를 중앙에
        run = tf.paragraphs[0].add_run()
        run.text = item
        apply_font(run, "body_a_bold" if version == "A" else "body_b", color=DARK_NAVY)


def build_text_heavy_table(slide, sc, version):
    """텍스트 + 테이블 (Version B)."""
    top = CONTENT_TOP

    # 본문 텍스트
    if sc.body_texts:
        text_height = min(Inches(2.5), CONTENT_HEIGHT * 0.4)
        _add_body_paragraphs(
            slide, CONTENT_LEFT, top,
            CONTENT_WIDTH, text_height,
            sc.body_texts[:5], "body_b", color=VERY_DARK_NAVY,
        )
        top += text_height + Inches(0.2)

    # 테이블
    if sc.tables:
        remaining = CONTENT_HEIGHT - (top - CONTENT_TOP)
        _add_styled_table(
            slide, CONTENT_LEFT, top,
            CONTENT_WIDTH, remaining, sc.tables[0],
        )


def build_thank_you(slide, sc, version):
    """감사 슬라이드 — 모던 디자인 (크롬 없음)."""
    # 배경 — 네이비 풀 블리드
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_NAVY
    bg.line.fill.background()

    # 좌측 Cyan 액센트 스트립 (표지와 통일)
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        Inches(0.12), SLIDE_HEIGHT,
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = BRAND_CYAN
    accent.line.fill.background()

    # 감사합니다 (좌측 정렬, 큰 타이포)
    _add_text_box(
        slide, Inches(2.0), Inches(1.8), Inches(9.0), Inches(1.2),
        "감사합니다", "main_heading", override_size=Pt(52),
        color=WHITE, alignment=PP_ALIGN.LEFT,
    )

    # 수평선
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2.0), Inches(3.1),
        Inches(3.0), Inches(0.03),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BRAND_CYAN
    line.line.fill.background()

    _add_text_box(
        slide, Inches(2.0), Inches(3.3), Inches(9.0), Inches(0.6),
        "Thank you for your time",
        "section_header", color=BRAND_CYAN, alignment=PP_ALIGN.LEFT,
    )

    # 연락처 영역 — 2컬럼
    _add_text_box(
        slide, Inches(2.0), Inches(4.5), Inches(4.0), Inches(0.35),
        "Contact", "body_a_bold", color=BRIGHT_BLUE, alignment=PP_ALIGN.LEFT,
    )
    contacts_left = [
        "sale@aisirius.ai",
        "031-360-7869",
        "www.aisirius.ai",
    ]
    c_top = Inches(5.0)
    for text in contacts_left:
        _add_text_box(
            slide, Inches(2.0), c_top, Inches(4.0), Inches(0.3),
            text, "small_label", color=WHITE, alignment=PP_ALIGN.LEFT,
        )
        c_top += Inches(0.35)

    _add_text_box(
        slide, Inches(7.0), Inches(4.5), Inches(5.0), Inches(0.35),
        "Address", "body_a_bold", color=BRIGHT_BLUE, alignment=PP_ALIGN.LEFT,
    )
    _add_text_box(
        slide, Inches(7.0), Inches(5.0), Inches(5.0), Inches(0.3),
        "경기도 수원시 영통구 광교로 105", "small_label",
        color=WHITE, alignment=PP_ALIGN.LEFT,
    )
    _add_text_box(
        slide, Inches(7.0), Inches(5.35), Inches(5.0), Inches(0.3),
        "경기R&DB센터 / 경기창조경제혁신센터 709호", "small_label",
        color=WHITE, alignment=PP_ALIGN.LEFT,
    )

    # CMS 데모
    _add_text_box(
        slide, Inches(7.0), Inches(5.85), Inches(5.0), Inches(0.3),
        "CMS Demo: cms.aisirius.ai", "small_label",
        color=MUTED_TEAL, alignment=PP_ALIGN.LEFT,
    )

    # 하단 3색 액센트 라인
    bar_h = Inches(0.06)
    bar_top = SLIDE_HEIGHT - bar_h
    bar_w = SLIDE_WIDTH / 3
    for i, color in enumerate([DARK_NAVY, STANDARD_BLUE, BRAND_CYAN]):
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(i * 13.333 / 3), bar_top, bar_w, bar_h,
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()


def build_generic(slide, sc, version):
    """범용 레이아웃 (매칭 안 되는 경우)."""
    if version == "B":
        build_text_heavy(slide, sc, version)
    else:
        build_two_column(slide, sc, version)


# ── 레이아웃 디스패처 ─────────────────────────────────
LAYOUT_BUILDERS = {
    LayoutType.TITLE_COVER: build_title_cover,
    LayoutType.TWO_COLUMN: build_two_column,
    LayoutType.CONTENT_IMAGE_LEFT: build_two_column,
    LayoutType.CONTENT_IMAGE_RIGHT: build_two_column,
    LayoutType.FULL_BLEED_IMAGE: build_two_column,  # 이미지 없으면 two-column으로 대체
    LayoutType.CHART_SLIDE: build_chart_slide,
    LayoutType.CHART_DETAILED: build_chart_detailed,
    LayoutType.COMPARISON_TABLE: build_comparison_table,
    LayoutType.THREE_PILLAR: build_three_pillar,
    LayoutType.INFOGRAPHIC_NUMBERS: build_infographic_numbers,
    LayoutType.TEXT_HEAVY_BULLETS: build_text_heavy,
    LayoutType.TEXT_HEAVY_TABLE: build_text_heavy_table,
    LayoutType.TEXT_SMALL_IMAGE: build_two_column,
    LayoutType.NEWS_QUOTE: build_news_quote,
    LayoutType.TWO_COLUMN_TEXT: build_text_heavy,
    LayoutType.TIMELINE: build_two_column,
    LayoutType.THANK_YOU: build_thank_you,
    LayoutType.SECTION_HEADER: build_text_heavy,
    LayoutType.GENERIC: build_generic,
}


# ══════════════════════════════════════════════════════
#  메인 빌드 함수
# ══════════════════════════════════════════════════════

def build_deck(version: str) -> str:
    """PPTX 빌드 실행.

    Args:
        version: "A" 또는 "B"

    Returns:
        출력 파일 경로
    """
    # 템플릿 로드
    template_path = os.path.join(TEMPLATES_DIR, f"aisirius_template_{version}.pptx")
    if not os.path.exists(template_path):
        print(f"[WARNING] 템플릿 없음: {template_path}, 빈 프레젠테이션으로 생성")
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
    else:
        prs = Presentation(template_path)

    # 콘텐츠 파싱
    if version == "A":
        content_path = os.path.join(CONTENT_DIR, "슬라이드_세부내용_VerA_비주얼.md")
    else:
        content_path = os.path.join(CONTENT_DIR, "슬라이드_세부내용_VerB_텍스트상세.md")

    slides_content = parse_slides(content_path, version)
    print(f"[INFO] {len(slides_content)}장 슬라이드 파싱 완료: {os.path.basename(content_path)}")

    section_map = SECTION_MAP_A if version == "A" else SECTION_MAP_B
    total_slides = len(slides_content)

    # 슬라이드별 빌드
    for sc in slides_content:
        # Blank 레이아웃 추가 (인덱스 6)
        layout_idx = min(6, len(prs.slide_layouts) - 1)
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # 크롬 적용 (표지/감사 제외)
        is_no_chrome = sc.layout_type in (LayoutType.TITLE_COVER, LayoutType.THANK_YOU)
        if not is_no_chrome:
            section_title, section_english = section_map.get(
                sc.number, ("", "")
            )
            add_full_chrome(
                slide,
                section_title=section_title,
                section_english=section_english,
                slide_number=sc.number,
                total_slides=total_slides,
            )

        # 레이아웃 빌더 실행
        builder_fn = LAYOUT_BUILDERS.get(sc.layout_type, build_generic)
        try:
            builder_fn(slide, sc, version)
        except Exception as e:
            print(f"[ERROR] Slide {sc.number} ({sc.title}): {e}")
            # 에러 시 제목이라도 표시
            _add_text_box(
                slide, CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, Inches(1),
                f"Slide {sc.number}: {sc.title}\n[빌드 오류: {e}]",
                "body_a", color=RGBColor(0xFF, 0x00, 0x00),
            )

        print(f"  [{sc.number:2d}/{total_slides}] {sc.title} ({sc.layout_type.value})")

    # 저장 (버전 관리: 기존 파일 덮어쓰지 않고 날짜+순번 부여)
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    today = date.today().strftime("%Y%m%d")
    prefix = f"AIsirius_회사소개서_Ver{version}_{today}"

    # 같은 날짜의 기존 파일 중 최대 순번 찾기
    existing = [
        f for f in os.listdir(OUTPUT_DIR)
        if f.startswith(prefix) and f.endswith(".pptx")
    ]
    max_seq = 0
    for f in existing:
        match = re.search(r"_(\d{2})\.pptx$", f)
        if match:
            max_seq = max(max_seq, int(match.group(1)))

    next_seq = max_seq + 1
    output_filename = f"{prefix}_{next_seq:02d}.pptx"
    output_path = os.path.join(OUTPUT_DIR, output_filename)
    prs.save(output_path)
    print(f"\n[OK] 저장 완료: {output_path}")
    print(f"     (기존 파일 {len(existing)}개 보존됨)")
    return output_path


def main():
    parser = argparse.ArgumentParser(description="AIsirius PPTX Builder")
    parser.add_argument("--version", choices=["A", "B", "all"], default="all",
                        help="빌드 버전 (A=비주얼, B=텍스트상세, all=둘 다)")
    parser.add_argument("--watermark", action="store_true",
                        help="워터마크 + PDF 변환 실행")
    parser.add_argument("--company", default="",
                        help="워터마크 대상 회사명")
    parser.add_argument("--password", default="",
                        help="PDF 암호")
    args = parser.parse_args()

    print("=" * 60)
    print("  AIsirius PPTX Builder")
    print("=" * 60)

    versions = ["A", "B"] if args.version == "all" else [args.version]

    for ver in versions:
        print(f"\n{'─' * 50}")
        print(f"  Version {ver} 빌드 시작")
        print(f"{'─' * 50}")
        output_path = build_deck(ver)

        # 워터마크 옵션
        if args.watermark and args.company:
            try:
                from watermark_pdf import generate_watermark_pdf
                logo_path = os.path.join(ASSETS_DIR, "icons", "aisirius_logo.png")
                if not os.path.exists(logo_path):
                    print(f"[WARNING] 로고 파일 없음: {logo_path}, 워터마크 건너뜀")
                else:
                    pdf_path, wm_path = generate_watermark_pdf(
                        input_pptx=output_path,
                        logo_path=logo_path,
                        company_name=args.company,
                        password=args.password or None,
                    )
                    print(f"[OK] 워터마크 PPTX: {wm_path}")
                    print(f"[OK] PDF: {pdf_path}")
            except ImportError:
                print("[WARNING] watermark_pdf 모듈 로드 실패, 워터마크 건너뜀")

    print(f"\n{'=' * 60}")
    print("  빌드 완료!")
    print(f"{'=' * 60}")


if __name__ == "__main__":
    main()
