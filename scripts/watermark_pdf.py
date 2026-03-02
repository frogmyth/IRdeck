"""
AIsirius PPTX → 워터마크 PDF 자동 생성 스크립트

사용법:
    python watermark_pdf.py <input.pptx> --logo <logo.png> [--company "회사명"] [--password "비밀번호"] [--output <output.pdf>]

프로세스:
    1. 회사 로고 PNG → 반투명 워터마크 이미지 자동 생성
    2. PPTX 각 슬라이드에 워터마크 삽입
    3. 워터마크 삽입된 PPTX를 output 폴더에 저장 (업체별 커스텀 가능)
    4. PPTX → PDF 변환 (LibreOffice soffice CLI)
    5. PDF에 암호 설정 (선택)

필요 라이브러리:
    pip install python-pptx Pillow pikepdf

외부 의존:
    - LibreOffice (PPTX→PDF 변환용) 또는 Windows의 경우 comtypes + PowerPoint
"""

import argparse
import os
import sys
import shutil
import subprocess
import tempfile
from pathlib import Path

from PIL import Image, ImageEnhance, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE


def create_watermark_image(
    logo_path: str,
    output_path: str,
    target_width: int = 800,
    opacity: float = 0.15,
    company_name: str = None,
    tile: bool = False,
    angle: float = -30,
):
    """회사 로고로부터 반투명 워터마크 이미지를 자동 생성한다.

    Args:
        logo_path: 원본 회사 로고 PNG 경로
        output_path: 생성할 워터마크 이미지 경로
        target_width: 워터마크 너비 (px)
        opacity: 투명도 (0.0~1.0, 낮을수록 투명)
        company_name: 로고 아래에 표시할 회사명 (None이면 로고만)
        tile: True면 타일 패턴 워터마크 생성
        angle: 회전 각도
    """
    logo = Image.open(logo_path).convert("RGBA")

    # 비율 유지하며 리사이즈
    ratio = target_width / logo.width
    new_height = int(logo.height * ratio)
    logo = logo.resize((target_width, new_height), Image.LANCZOS)

    # 회사명 추가 (옵션)
    if company_name:
        try:
            font = ImageFont.truetype("malgun.ttf", size=int(target_width * 0.06))
        except OSError:
            font = ImageFont.load_default()

        text_layer = Image.new("RGBA", (target_width, new_height + 60), (0, 0, 0, 0))
        text_layer.paste(logo, (0, 0))
        draw = ImageDraw.Draw(text_layer)
        bbox = draw.textbbox((0, 0), company_name, font=font)
        text_w = bbox[2] - bbox[0]
        x = (target_width - text_w) // 2
        draw.text((x, new_height + 5), company_name, fill=(128, 128, 128, 255), font=font)
        logo = text_layer
        new_height += 60

    # 투명도 적용
    r, g, b, a = logo.split()
    a = ImageEnhance.Brightness(a).enhance(opacity)
    logo = Image.merge("RGBA", (r, g, b, a))

    # 회전
    if angle != 0:
        logo = logo.rotate(angle, expand=True, resample=Image.BICUBIC)

    if tile:
        # 타일 패턴: 슬라이드 크기(16:9 = 1920x1080)에 맞춰 반복
        canvas_w, canvas_h = 1920, 1080
        tile_w, tile_h = logo.size
        canvas = Image.new("RGBA", (canvas_w, canvas_h), (0, 0, 0, 0))
        spacing_x = tile_w + 100
        spacing_y = tile_h + 100
        for y in range(-tile_h, canvas_h + tile_h, spacing_y):
            for x in range(-tile_w, canvas_w + tile_w, spacing_x):
                canvas.paste(logo, (x, y), logo)
        canvas.save(output_path, "PNG")
    else:
        logo.save(output_path, "PNG")

    return output_path


def add_watermark_to_pptx(
    input_pptx: str,
    output_pptx: str,
    watermark_path: str,
    position: str = "center",
):
    """PPTX 파일의 모든 슬라이드에 워터마크 이미지를 삽입한다.

    Args:
        input_pptx: 원본 PPTX 경로
        output_pptx: 워터마크 삽입된 PPTX 경로
        watermark_path: 워터마크 이미지 경로
        position: 'center', 'bottom-right', 'tile'
    """
    prs = Presentation(input_pptx)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    wm_img = Image.open(watermark_path)
    wm_width, wm_height = wm_img.size

    for slide in prs.slides:
        if position == "center":
            # 슬라이드 중앙에 워터마크 배치
            scale = min(
                (slide_width * 0.6) / Emu(wm_width * 914400 / 96),
                (slide_height * 0.6) / Emu(wm_height * 914400 / 96),
            )
            img_width = Emu(int(wm_width * 914400 / 96 * scale))
            img_height = Emu(int(wm_height * 914400 / 96 * scale))
            left = (slide_width - img_width) // 2
            top = (slide_height - img_height) // 2

        elif position == "bottom-right":
            scale = 0.2
            img_width = Emu(int(wm_width * 914400 / 96 * scale))
            img_height = Emu(int(wm_height * 914400 / 96 * scale))
            left = slide_width - img_width - Inches(0.3)
            top = slide_height - img_height - Inches(0.3)

        elif position == "tile":
            # 타일의 경우 전체 슬라이드 크기로 배치
            img_width = slide_width
            img_height = slide_height
            left = 0
            top = 0

        pic = slide.shapes.add_picture(watermark_path, left, top, img_width, img_height)

        # 워터마크를 맨 뒤로 보내기 (다른 콘텐츠 위에 겹치지 않도록)
        sp = pic._element
        sp.getparent().remove(sp)
        slide.shapes._spTree.insert(2, sp)  # 배경 바로 위

    prs.save(output_pptx)
    return output_pptx


def pptx_to_pdf(input_pptx: str, output_pdf: str):
    """PPTX를 PDF로 변환한다.

    Windows: PowerPoint COM 자동화 (comtypes) 시도, 없으면 LibreOffice
    Linux/Mac: LibreOffice soffice CLI
    """
    output_dir = str(Path(output_pdf).parent)

    if sys.platform == "win32":
        # Method 1: PowerPoint COM (가장 정확)
        try:
            import comtypes.client

            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            powerpoint.Visible = 1
            abs_pptx = os.path.abspath(input_pptx)
            abs_pdf = os.path.abspath(output_pdf)
            deck = powerpoint.Presentations.Open(abs_pptx, WithWindow=False)
            deck.SaveAs(abs_pdf, 32)  # 32 = ppSaveAsPDF
            deck.Close()
            powerpoint.Quit()
            return output_pdf
        except Exception:
            pass  # PowerPoint not installed, try LibreOffice

    # Method 2: LibreOffice CLI
    soffice_paths = [
        "soffice",
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        "/usr/bin/soffice",
        "/usr/bin/libreoffice",
    ]

    for soffice in soffice_paths:
        try:
            result = subprocess.run(
                [soffice, "--headless", "--convert-to", "pdf", "--outdir", output_dir, input_pptx],
                capture_output=True,
                text=True,
                timeout=120,
            )
            if result.returncode == 0:
                # LibreOffice는 원본 파일명.pdf로 생성
                generated = Path(output_dir) / (Path(input_pptx).stem + ".pdf")
                if str(generated) != output_pdf and generated.exists():
                    shutil.move(str(generated), output_pdf)
                return output_pdf
        except (FileNotFoundError, subprocess.TimeoutExpired):
            continue

    raise RuntimeError(
        "PDF 변환 실패: PowerPoint 또는 LibreOffice가 필요합니다.\n"
        "- Windows: Microsoft PowerPoint 설치 또는 LibreOffice 설치\n"
        "- Linux/Mac: sudo apt install libreoffice (또는 brew install libreoffice)"
    )


def protect_pdf(input_pdf: str, output_pdf: str, password: str):
    """PDF에 암호를 설정한다.

    Args:
        input_pdf: 원본 PDF 경로
        output_pdf: 암호 설정된 PDF 경로
        password: 열기 암호
    """
    import pikepdf

    pdf = pikepdf.open(input_pdf)
    pdf.save(
        output_pdf,
        encryption=pikepdf.Encryption(
            owner=password + "_owner",  # 소유자 암호 (편집/인쇄 제한)
            user=password,  # 열기 암호
            R=6,  # AES-256 암호화
            allow=pikepdf.Permissions(
                extract=False,       # 텍스트/이미지 추출 불가
                modify_annotation=False,
                modify_other=False,
                modify_assembly=False,
                print_lowres=True,   # 저해상도 인쇄만 허용
                print_highres=False, # 고해상도 인쇄 불가
            ),
        ),
    )
    pdf.close()
    return output_pdf


def generate_watermark_pdf(
    input_pptx: str,
    logo_path: str,
    company_name: str = None,
    password: str = None,
    output_pdf: str = None,
    watermark_style: str = "center",
    opacity: float = 0.15,
):
    """전체 프로세스: PPTX → 워터마크 삽입 → PDF → 암호 설정

    Args:
        input_pptx: 원본 PPTX 파일 경로
        logo_path: 회사 로고 PNG 경로
        company_name: 워터마크에 표시할 회사명 (옵션)
        password: PDF 열기 암호 (옵션, None이면 암호 미설정)
        output_pdf: 출력 PDF 경로 (None이면 자동 생성)
        watermark_style: 'center', 'bottom-right', 'tile'
        opacity: 워터마크 투명도 (0.0~1.0)

    Returns:
        tuple: (생성된 PDF 파일 경로, 워터마크 삽입된 PPTX 파일 경로)
    """
    stem = Path(input_pptx).stem
    output_dir = Path(input_pptx).parent

    if output_pdf is None:
        suffix = "+watermark"
        if password:
            suffix += "+protected"
        output_pdf = str(output_dir / f"{stem}{suffix}.pdf")

    # 워터마크 PPTX는 항상 저장 (업체별 내용 커스텀 가능)
    output_pptx_wm = str(output_dir / f"{stem}+watermark.pptx")

    with tempfile.TemporaryDirectory() as tmpdir:
        # Step 1: 로고 → 워터마크 이미지 생성
        print(f"[1/4] 워터마크 이미지 생성 중... (스타일: {watermark_style}, 투명도: {opacity})")
        tile = watermark_style == "tile"
        wm_path = os.path.join(tmpdir, "watermark.png")
        create_watermark_image(
            logo_path=logo_path,
            output_path=wm_path,
            opacity=opacity,
            company_name=company_name,
            tile=tile,
        )

        # Step 2: PPTX에 워터마크 삽입 → 영구 저장
        print(f"[2/4] PPTX에 워터마크 삽입 중... → {output_pptx_wm}")
        add_watermark_to_pptx(
            input_pptx=input_pptx,
            output_pptx=output_pptx_wm,
            watermark_path=wm_path,
            position=watermark_style,
        )

        # Step 3: PDF 변환
        print("[3/4] PDF 변환 중...")
        if password:
            tmp_pdf = os.path.join(tmpdir, "temp.pdf")
            pptx_to_pdf(output_pptx_wm, tmp_pdf)

            # Step 4: 암호 설정
            print("[4/4] PDF 암호 설정 중... (AES-256)")
            protect_pdf(tmp_pdf, output_pdf, password)
        else:
            pptx_to_pdf(output_pptx_wm, output_pdf)
            print("[4/4] 암호 설정 건너뜀 (--password 미지정)")

    print(f"\n완료:")
    print(f"  PPTX (워터마크): {output_pptx_wm}")
    print(f"  PDF: {output_pdf}")
    print(f"  ※ 워터마크 PPTX는 업체별 내용 커스텀 후 재변환 가능")
    return output_pdf, output_pptx_wm


def main():
    parser = argparse.ArgumentParser(
        description="AIsirius PPTX → 워터마크 PDF 자동 생성",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
사용 예시:
  # 기본 (중앙 워터마크, 암호 없음)
  python watermark_pdf.py presentation.pptx --logo logo.png

  # 회사명 포함 + 암호 설정
  python watermark_pdf.py presentation.pptx --logo logo.png --company "AIsirius" --password "secret123"

  # 타일 패턴 워터마크 + 고투명도
  python watermark_pdf.py presentation.pptx --logo logo.png --style tile --opacity 0.08

  # 출력 파일 지정
  python watermark_pdf.py presentation.pptx --logo logo.png -o output/final.pdf
""",
    )
    parser.add_argument("input_pptx", help="원본 PPTX 파일 경로")
    parser.add_argument("--logo", required=True, help="회사 로고 PNG 파일 경로")
    parser.add_argument("--company", default=None, help="워터마크에 표시할 회사명")
    parser.add_argument("--password", default=None, help="PDF 열기 암호 (미지정 시 암호 없음)")
    parser.add_argument("-o", "--output", default=None, help="출력 PDF 파일 경로")
    parser.add_argument(
        "--style",
        choices=["center", "bottom-right", "tile"],
        default="center",
        help="워터마크 스타일 (기본: center)",
    )
    parser.add_argument(
        "--opacity",
        type=float,
        default=0.15,
        help="워터마크 투명도 0.0~1.0 (기본: 0.15)",
    )

    args = parser.parse_args()

    if not os.path.exists(args.input_pptx):
        print(f"오류: PPTX 파일을 찾을 수 없습니다: {args.input_pptx}")
        sys.exit(1)
    if not os.path.exists(args.logo):
        print(f"오류: 로고 파일을 찾을 수 없습니다: {args.logo}")
        sys.exit(1)

    generate_watermark_pdf(
        input_pptx=args.input_pptx,
        logo_path=args.logo,
        company_name=args.company,
        password=args.password,
        output_pdf=args.output,
        watermark_style=args.style,
        opacity=args.opacity,
    )


if __name__ == "__main__":
    main()
