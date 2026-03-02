"""슬라이드 크롬 요소: 헤더바, 푸터, 로고, 분리선, 장식 요소.

모던 미니멀 디자인 적용:
- 슬림 헤더 + 좌측 Cyan 액센트 바
- 하단 슬림 액센트 라인 (3색 → 그라데이션 효과)
- 우상단 로고
- 좌하단 섹션 인디케이터
"""
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

from .colors import (
    DARK_NAVY, BRAND_CYAN, STANDARD_BLUE, WHITE, MUTED_TEAL,
)
from .fonts import apply_font
from .layouts import (
    SLIDE_WIDTH, SLIDE_HEIGHT,
    HEADER_BAR_TOP,
    LOGO_LEFT, LOGO_WIDTH,
)


def add_header_bar(slide, section_title="", section_english=""):
    """모던 헤더: 좌측 Cyan 액센트 스트립 + 상단 네이비 바."""
    # 좌측 세로 Cyan 액센트 스트립 (상단 ~ 콘텐츠 영역까지)
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(0.12), SLIDE_HEIGHT,
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = BRAND_CYAN
    accent.line.fill.background()

    # 상단 네이비 헤더 바 (슬림)
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.12), HEADER_BAR_TOP,
        SLIDE_WIDTH - Inches(0.12), Inches(0.65),
    )
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_NAVY
    header.line.fill.background()

    # 섹션 제목 텍스트
    if section_title or section_english:
        txbox = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.08),
            Inches(8.0), Inches(0.55),
        )
        tf = txbox.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)

        # 영문 라벨 (작게)
        if section_english:
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = section_english
            apply_font(run, "header_english", color=BRAND_CYAN)

        # 한글 제목 (같은 줄, 구분자 + 큰 텍스트)
        if section_title:
            if section_english:
                p = tf.paragraphs[0]
                sep = p.add_run()
                sep.text = "  |  "
                apply_font(sep, "header_english", color=MUTED_TEAL)
                run = p.add_run()
                run.text = section_title
                apply_font(run, "header_english", override_size=Pt(16), override_bold=True, color=WHITE)
            else:
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = section_title
                apply_font(run, "section_header", override_size=Pt(18), color=WHITE)


def add_footer_bars(slide):
    """모던 푸터: 슬림 3색 액센트 라인 (하단 4px)."""
    bar_height = Inches(0.06)  # 약 4px
    bar_top = SLIDE_HEIGHT - bar_height
    bar_width = SLIDE_WIDTH / 3

    colors = [DARK_NAVY, STANDARD_BLUE, BRAND_CYAN]
    for i, color in enumerate(colors):
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(i * 13.333 / 3), bar_top,
            bar_width, bar_height,
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()


def add_separator_line(slide):
    """헤더 아래 Cyan 분리선."""
    connector = slide.shapes.add_connector(
        1,  # straight connector
        Inches(0.12), Inches(0.66),
        SLIDE_WIDTH, Inches(0.66),
    )
    connector.line.color.rgb = BRAND_CYAN
    connector.line.width = Pt(1.0)


def add_logo(slide):
    """우상단 AIsirius 로고 텍스트."""
    txbox = slide.shapes.add_textbox(
        LOGO_LEFT, Inches(0.12),
        LOGO_WIDTH, Inches(0.42),
    )
    tf = txbox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT

    # "AI" run
    run_ai = p.add_run()
    run_ai.text = "AI"
    apply_font(run_ai, "logo_ai", color=BRAND_CYAN)

    # "sirius" run
    run_sirius = p.add_run()
    run_sirius.text = "sirius"
    apply_font(run_sirius, "logo_sirius", color=WHITE)


def add_slide_number(slide, number, total=None):
    """슬라이드 번호 (우하단, 푸터 위)."""
    txbox = slide.shapes.add_textbox(
        Inches(12.2), SLIDE_HEIGHT - Inches(0.4),
        Inches(0.9), Inches(0.3),
    )
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT

    text = f"{number}" if total is None else f"{number} / {total}"
    run = p.add_run()
    run.text = text
    apply_font(run, "slide_number", color=MUTED_TEAL)


def add_bottom_accent(slide):
    """좌하단 장식 요소: 작은 Cyan 사각형 + 회사명."""
    # 작은 Cyan 사각형 (인디케이터)
    indicator = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), SLIDE_HEIGHT - Inches(0.45),
        Inches(0.3), Inches(0.08),
    )
    indicator.fill.solid()
    indicator.fill.fore_color.rgb = BRAND_CYAN
    indicator.line.fill.background()

    # 회사명
    txbox = slide.shapes.add_textbox(
        Inches(0.9), SLIDE_HEIGHT - Inches(0.5),
        Inches(3.0), Inches(0.3),
    )
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "AIsirius Co., Ltd.  |  Confidential"
    apply_font(run, "source_text", color=MUTED_TEAL)


def add_full_chrome(slide, section_title="", section_english="", slide_number=None, total_slides=None):
    """크롬 요소 전체 적용 (표지/감사 슬라이드 제외).

    Args:
        slide: pptx Slide 객체
        section_title: 한글 섹션 제목
        section_english: 영문 라벨
        slide_number: 슬라이드 번호 (int, optional)
        total_slides: 전체 슬라이드 수 (int, optional)
    """
    add_header_bar(slide, section_title, section_english)
    add_separator_line(slide)
    add_footer_bars(slide)
    add_logo(slide)
    add_bottom_accent(slide)
    if slide_number is not None:
        add_slide_number(slide, slide_number, total_slides)
