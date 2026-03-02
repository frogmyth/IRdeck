"""시장 데이터 차트 생성 (python-pptx 네이티브 차트).

market_data.md의 데이터를 사용하여 6종의 차트를 생성한다.
PowerPoint에서 편집 가능한 네이티브 차트 객체.
"""
from pptx.util import Inches, Pt, Emu
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.dml.color import RGBColor

from .colors import (
    DARK_NAVY, STANDARD_BLUE, BRAND_CYAN, TEAL, LIME_GREEN,
    BRIGHT_BLUE, MEDIUM_BLUE, VERY_DARK_NAVY, WHITE, LIGHT_GRAY,
    MUTED_TEAL, CHART_SERIES,
)
from .fonts import apply_font


# ── 시장 데이터 (market_data.md 기반) ─────────────────
MARKET_DATA = {
    "ESL": {
        "years": [2022, 2023, 2024, 2025, 2026, 2027, 2028, 2029, 2030],
        "values": [1.29, 1.61, 2.34, 2.75, 2.39, 3.10, 3.53, 4.18, 3.78],
        "label": "글로벌 ESL 시장 ($B)",
        "cagr": "13.9%",
        "source": "MarketsandMarkets / Mordor Intelligence",
    },
    "RMN_AD_SPEND": {
        "years": [2022, 2023, 2024, 2025, 2026, 2027, 2028, 2029, 2030],
        "values": [114.6, 128.2, 141.7, 184.0, 196.7, 218.4, 242.4, 269.1, 312.0],
        "label": "글로벌 RMN 광고비 ($B)",
        "cagr": "11.0%",
        "source": "Forrester / WARC Media",
    },
    "AI_IN_RETAIL": {
        "years": [2022, 2023, 2024, 2025, 2026, 2027, 2028, 2029, 2030],
        "values": [5.50, 9.36, 11.61, 14.49, 17.82, 21.92, 26.96, 33.16, 40.74],
        "label": "글로벌 AI in Retail ($B)",
        "cagr": "23.0%",
        "source": "Grand View Research",
    },
    "WALMART_CONNECT": {
        "years": ["FY2021", "FY2022", "FY2023", "FY2024", "FY2025"],
        "values": [2.10, 2.70, 3.40, 4.40, 6.40],
        "label": "Walmart Connect 광고 매출 ($B)",
        "source": "Adweek / Statista / AdExchanger",
    },
    "J_CURVE": {
        "years": [2025, 2026, 2027, 2028, 2029, 2030],
        "values": [10, 35, 92, 247, 550, 1100],
        "label": "AIsirius 매출 로드맵 (억원)",
        "milestones": ["Pre-A", "35억", "시리즈A", "247억", "시리즈B", "1,100억/IPO"],
        "source": "AIsirius 내부 계획",
    },
}

# ESL 지역별 도입률 데이터
ESL_ADOPTION = {
    "headers": ["지역", "유통사 수", "ESL 도입률", "교체 시점"],
    "rows": [
        ["미국", "77", "5%", "신규 진입"],
        ["일본", "29", "15%", "2027~2028"],
        ["유럽", "63", "30%", "2027~2028"],
        ["한국", "5", "80%", "2027~2028"],
        ["기타", "76", "<1%", "미개척"],
    ],
}


def _style_chart_common(chart):
    """차트 공통 스타일 적용."""
    chart.has_legend = False

    # 값 축
    value_axis = chart.value_axis
    value_axis.has_title = False
    value_axis.major_gridlines.format.line.color.rgb = LIGHT_GRAY
    value_axis.major_gridlines.format.line.width = Pt(0.5)
    value_axis.format.line.color.rgb = MUTED_TEAL
    value_axis.format.line.width = Pt(0.5)
    value_axis.tick_labels.font.size = Pt(9)
    value_axis.tick_labels.font.color.rgb = VERY_DARK_NAVY

    # 카테고리 축
    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(9)
    category_axis.tick_labels.font.color.rgb = VERY_DARK_NAVY
    category_axis.format.line.color.rgb = MUTED_TEAL
    category_axis.format.line.width = Pt(0.5)


def _color_bar_series(plot, color):
    """막대 차트 시리즈 색상 설정."""
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = color


def _color_line_series(plot, color, width=Pt(2.5)):
    """라인 차트 시리즈 스타일 설정."""
    series = plot.series[0]
    series.format.line.color.rgb = color
    series.format.line.width = width
    series.smooth = True


def add_source_text(slide, left, top, width, source_text):
    """차트 하단에 출처 텍스트 추가."""
    txbox = slide.shapes.add_textbox(left, top, width, Inches(0.3))
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"출처: {source_text}"
    apply_font(run, "source_text", color=MUTED_TEAL)


def add_esl_chart(slide, left, top, width, height):
    """ESL 시장 막대 차트 (2022-2030)."""
    data = MARKET_DATA["ESL"]
    chart_data = CategoryChartData()
    chart_data.categories = [str(y) for y in data["years"]]
    chart_data.add_series(data["label"], data["values"])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    _style_chart_common(chart)
    _color_bar_series(chart.plots[0], DARK_NAVY)

    # 타이틀
    chart.has_title = True
    chart.chart_title.text_frame.paragraphs[0].text = f"ESL 시장 (CAGR {data['cagr']})"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = DARK_NAVY
    chart.chart_title.text_frame.paragraphs[0].font.bold = True

    # 출처
    add_source_text(slide, left, top + height, width, data["source"])
    return chart_frame


def add_rmn_chart(slide, left, top, width, height):
    """RMN 광고비 라인 차트 (2022-2030)."""
    data = MARKET_DATA["RMN_AD_SPEND"]
    chart_data = CategoryChartData()
    chart_data.categories = [str(y) for y in data["years"]]
    chart_data.add_series(data["label"], data["values"])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    _style_chart_common(chart)
    _color_line_series(chart.plots[0], STANDARD_BLUE)

    chart.has_title = True
    chart.chart_title.text_frame.paragraphs[0].text = f"RMN 광고비 (CAGR {data['cagr']})"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = DARK_NAVY
    chart.chart_title.text_frame.paragraphs[0].font.bold = True

    add_source_text(slide, left, top + height, width, data["source"])
    return chart_frame


def add_ai_retail_chart(slide, left, top, width, height):
    """AI in Retail 막대 차트 (2022-2030)."""
    data = MARKET_DATA["AI_IN_RETAIL"]
    chart_data = CategoryChartData()
    chart_data.categories = [str(y) for y in data["years"]]
    chart_data.add_series(data["label"], data["values"])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    _style_chart_common(chart)
    _color_bar_series(chart.plots[0], TEAL)

    chart.has_title = True
    chart.chart_title.text_frame.paragraphs[0].text = f"AI in Retail (CAGR {data['cagr']})"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = DARK_NAVY
    chart.chart_title.text_frame.paragraphs[0].font.bold = True

    add_source_text(slide, left, top + height, width, data["source"])
    return chart_frame


def add_walmart_chart(slide, left, top, width, height):
    """Walmart Connect 막대 차트 (FY2021-FY2025)."""
    data = MARKET_DATA["WALMART_CONNECT"]
    chart_data = CategoryChartData()
    chart_data.categories = data["years"]
    chart_data.add_series(data["label"], data["values"])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    _style_chart_common(chart)
    _color_bar_series(chart.plots[0], BRIGHT_BLUE)

    chart.has_title = True
    chart.chart_title.text_frame.paragraphs[0].text = "Walmart Connect 광고 매출"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = DARK_NAVY
    chart.chart_title.text_frame.paragraphs[0].font.bold = True

    add_source_text(slide, left, top + height, width, data["source"])
    return chart_frame


def add_jcurve_chart(slide, left, top, width, height):
    """J-Curve 매출 로드맵 라인 차트 (2025-2030)."""
    data = MARKET_DATA["J_CURVE"]
    chart_data = CategoryChartData()
    chart_data.categories = [str(y) for y in data["years"]]
    chart_data.add_series(data["label"], data["values"])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    _style_chart_common(chart)
    _color_line_series(chart.plots[0], BRAND_CYAN, width=Pt(3))

    chart.has_title = True
    chart.chart_title.text_frame.paragraphs[0].text = "AIsirius 매출 로드맵 (억원)"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = DARK_NAVY
    chart.chart_title.text_frame.paragraphs[0].font.bold = True

    # 데이터 레이블 표시
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.size = Pt(10)
    data_labels.font.bold = True
    data_labels.font.color.rgb = DARK_NAVY
    data_labels.number_format = "#,##0"

    add_source_text(slide, left, top + height, width, data["source"])
    return chart_frame


def add_esl_adoption_table(slide, left, top, width, height):
    """ESL 지역별 도입률 테이블 (스타일 테이블).

    python-pptx Table 객체로 생성.
    """
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    data = ESL_ADOPTION
    rows_count = len(data["rows"]) + 1  # +1 for header
    cols_count = len(data["headers"])

    table_shape = slide.shapes.add_table(
        rows_count, cols_count, left, top, width, height
    )
    table = table_shape.table

    # 열 너비 설정
    col_widths = [Inches(1.2), Inches(1.5), Inches(1.5), Inches(1.8)]
    for i, w in enumerate(col_widths[:cols_count]):
        table.columns[i].width = w

    # 헤더 행
    for j, header in enumerate(data["headers"]):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_NAVY
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.alignment = PP_ALIGN.CENTER

    # 데이터 행
    for i, row in enumerate(data["rows"]):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = val
            # 교대 행 색상
            if i % 2 == 0:
                from .colors import PALE_BLUE
                cell.fill.solid()
                cell.fill.fore_color.rgb = PALE_BLUE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.color.rgb = VERY_DARK_NAVY
                paragraph.alignment = PP_ALIGN.CENTER

    return table_shape


# ── 차트 키 → 함수 매핑 ──────────────────────────────
CHART_BUILDERS = {
    "ESL": add_esl_chart,
    "RMN_AD_SPEND": add_rmn_chart,
    "AI_IN_RETAIL": add_ai_retail_chart,
    "WALMART_CONNECT": add_walmart_chart,
    "J_CURVE": add_jcurve_chart,
}
