import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation(r"G:\00.googledrive\07.AIsirius\20.회사소개서\AIsirius_ 소개자료_신세계inc_20260226.pptx")

# Slide dimensions
print(f"=== Slide Dimensions ===")
print(f"Width: {prs.slide_width} EMU = {prs.slide_width / 914400:.2f} inches")
print(f"Height: {prs.slide_height} EMU = {prs.slide_height / 914400:.2f} inches")
print(f"Aspect: {'16:9' if abs(prs.slide_width/prs.slide_height - 16/9) < 0.01 else '4:3' if abs(prs.slide_width/prs.slide_height - 4/3) < 0.01 else 'custom'}")
print(f"Total slides: {len(prs.slides)}")

# Slide layouts
print(f"\n=== Slide Master & Layouts ===")
for i, master in enumerate(prs.slide_masters):
    print(f"Master {i}: '{master.name if hasattr(master, 'name') else 'unnamed'}'")
    for j, layout in enumerate(master.slide_layouts):
        print(f"  Layout {j}: {layout.name}")

# Analyze each slide (first 5)
print(f"\n=== Per-Slide Analysis (first 5 slides) ===")
slide_count = 0
for slide in prs.slides:
    slide_count += 1
    if slide_count > 5:
        break
    print(f"\n--- Slide {slide_count} ---")
    print(f"  Layout: {slide.slide_layout.name}")

    # Background
    try:
        bg = slide.background
        if bg.fill and bg.fill.type is not None:
            print(f"  Background fill type: {bg.fill.type}")
    except Exception as e:
        print(f"  Background: (could not read: {e})")

    for shape in slide.shapes:
        print(f"  Shape: {shape.shape_type}, name='{shape.name}', pos=({shape.left},{shape.top}), size=({shape.width},{shape.height})")

        if shape.has_text_frame:
            para_count = 0
            for para in shape.text_frame.paragraphs:
                para_count += 1
                if para_count > 3:
                    break
                align = para.alignment
                run_count = 0
                for run in para.runs:
                    run_count += 1
                    if run_count > 2:
                        break
                    font = run.font
                    text_display = f"'{run.text[:50]}...'" if len(run.text) > 50 else f"'{run.text}'"
                    print(f"    Text: {text_display}")
                    print(f"      Font: {font.name}, Size: {font.size}, Bold: {font.bold}, Italic: {font.italic}")
                    try:
                        if font.color and font.color.rgb:
                            print(f"      Color: #{font.color.rgb}")
                    except Exception:
                        pass

# Color analysis across all slides
print(f"\n=== Color Palette (all slides) ===")
colors = set()
fonts_used = set()
font_sizes = set()
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name:
                        fonts_used.add(run.font.name)
                    if run.font.size:
                        font_sizes.add(run.font.size)
                    try:
                        if run.font.color and run.font.color.rgb:
                            colors.add(str(run.font.color.rgb))
                    except Exception:
                        pass

print(f"Fonts: {sorted(fonts_used)}")
print(f"Font sizes (pt): {sorted([int(s/12700) for s in font_sizes])}")
print(f"Colors: {sorted(colors)}")

# Image analysis
print(f"\n=== Images ===")
img_count = 0
for slide_num_idx, slide in enumerate(prs.slides, 1):
    for shape in slide.shapes:
        if shape.shape_type == 13:  # Picture
            img_count += 1
            print(f"  Slide {slide_num_idx}: Image '{shape.name}', size=({shape.width/914400:.1f}x{shape.height/914400:.1f} in), pos=({shape.left/914400:.1f},{shape.top/914400:.1f} in)")
if img_count == 0:
    print("  No standalone images found.")

# Shape types summary
print(f"\n=== Shape Types Summary ===")
shape_types = {}
for slide in prs.slides:
    for shape in slide.shapes:
        st = str(shape.shape_type)
        shape_types[st] = shape_types.get(st, 0) + 1
for st, count in sorted(shape_types.items(), key=lambda x: -x[1]):
    print(f"  {st}: {count}")

# Remaining slides summary (6+)
print(f"\n=== Remaining Slides Layout Summary ===")
slide_idx = 0
for slide in prs.slides:
    slide_idx += 1
    if slide_idx <= 5:
        continue
    shape_count = sum(1 for _ in slide.shapes)
    text_shapes = sum(1 for s in slide.shapes if s.has_text_frame)
    print(f"  Slide {slide_idx}: Layout='{slide.slide_layout.name}', Shapes={shape_count}, TextShapes={text_shapes}")
